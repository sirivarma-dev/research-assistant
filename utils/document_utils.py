"""
Unified document text extraction.

Handles every upload type the app accepts:
  - PDF  (.pdf)   -> pypdf, with automatic OCR fallback for scanned PDFs
  - Word (.docx)  -> python-docx
  - Text (.txt/.md)
  - PowerPoint (.pptx) -> python-pptx
  - Images (.png/.jpg/.jpeg/.tiff/.bmp) -> OCR

Every function returns a (text, num_pages) tuple so the rest of the app can
treat all formats the same way.

OCR is OPTIONAL: it needs the free Tesseract program installed on the machine.
If Tesseract isn't available we degrade gracefully - normal documents still work
and scanned ones raise a clear, friendly error telling the user what to install.
"""
import os

from pypdf import PdfReader

# Image extensions we treat as "scanned" and send straight to OCR.
IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp")

# Everything the uploader is allowed to accept (used by app.py too).
SUPPORTED_EXTS = (".pdf", ".docx", ".txt", ".md", ".pptx") + IMAGE_EXTS


def ocr_available():
    """True only if the OCR stack (pytesseract + the Tesseract program) is ready."""
    try:
        import pytesseract  # noqa: F401
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


# ---------------------------------------------------------------- per-format --
def _extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return text.strip(), 1


def _extract_docx(path):
    import docx  # python-docx
    document = docx.Document(path)
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    # Also pull text out of any tables.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip(), 1


def _extract_pptx(path):
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n".join(parts).strip(), len(prs.slides)


def _ocr_image(path):
    """OCR a single image file."""
    import pytesseract
    from PIL import Image
    with Image.open(path) as img:
        return pytesseract.image_to_string(img).strip(), 1


def _ocr_pdf(path, dpi=200):
    """
    OCR a scanned PDF by rendering each page to an image with PyMuPDF (no
    Poppler needed) and running Tesseract on it.
    """
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image

    text_parts = []
    doc = fitz.open(path)
    zoom = dpi / 72  # 72 is the PDF's native DPI
    matrix = fitz.Matrix(zoom, zoom)
    try:
        for page in doc:
            pix = page.get_pixmap(matrix=matrix)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            text_parts.append(pytesseract.image_to_string(img))
        pages = doc.page_count
    finally:
        doc.close()
    return "\n".join(text_parts).strip(), pages


def _extract_pdf(path):
    """
    Normal (text-based) PDF: read with pypdf. If almost no text comes out, the
    PDF is probably scanned - fall back to OCR when it's available.
    """
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"
    text = text.strip()
    pages = len(reader.pages)

    # Heuristic: a real text PDF yields plenty of characters. Very little text
    # across multiple pages means it's a scan (images of text).
    looks_scanned = len(text) < 100 * max(pages, 1) and len(text) < 500
    if looks_scanned:
        if ocr_available():
            ocr_text, _ = _ocr_pdf(path)
            if ocr_text:
                return ocr_text, pages
        raise ValueError(
            "This looks like a scanned PDF (no selectable text). OCR is needed "
            "to read it, but Tesseract isn't installed. Install it from "
            "https://github.com/UB-Mannheim/tesseract/wiki and restart the app, "
            "or upload a text-based document instead."
        )
    return text, pages


# ------------------------------------------------------------------- public ---
def extract_text(path):
    """
    Extract (text, num_pages) from any supported document.
    Raises ValueError with a friendly message for unsupported types, empty
    documents, or scanned files when OCR isn't available.
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        text, pages = _extract_pdf(path)
    elif ext == ".docx":
        text, pages = _extract_docx(path)
    elif ext in (".txt", ".md"):
        text, pages = _extract_txt(path)
    elif ext == ".pptx":
        text, pages = _extract_pptx(path)
    elif ext in IMAGE_EXTS:
        if not ocr_available():
            raise ValueError(
                "Reading an image needs OCR, but Tesseract isn't installed. "
                "Install it from https://github.com/UB-Mannheim/tesseract/wiki "
                "and restart the app."
            )
        text, pages = _ocr_image(path)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. Please upload a PDF, DOCX, TXT, "
            "PPTX, or an image (PNG/JPG)."
        )

    if not text:
        raise ValueError(
            "No readable text was found in this document. It may be empty or, "
            "if scanned, require OCR."
        )
    return text, pages
